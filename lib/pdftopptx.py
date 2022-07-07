import os
import time
import shutil
import collections 
import collections.abc
import colorama
import subprocess

from PIL import Image
from pptx import Presentation
from tkinter import filedialog
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from colorama import Fore
from pdf2image import convert_from_path, convert_from_bytes
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)
colorama.init()

def cache_cleaner(directory): # Clear/Remove Cache
    cache_dir = directory + '/cache/Pdf to Pptx'
    shutil.rmtree(cache_dir)

def n_files(directory): # Count .pdf files in the directory
    total = 0

    for file in os.listdir(directory):
        if file.endswith('.pdf'):
            total += 1
    return total

# open file dialog using Tkinter GUI
def selectpdffile(directory): # Select and Copy file to directory
    filepath = filedialog.askopenfilename(initialdir="Documents",
                                          title="Select File",
                                          filetypes= (("pdf files","*.pdf"), # File Selection (Make sure user only input .pptx file)
                                          ))
    file = filepath
    file_dir = directory
    shutil.copy(file, file_dir) # Copy selected file into program directory

def createFolder(directory): # Create Output Folder
    try:
        if not os.path.exists(directory + 'Output Folder/Pdf to Pptx'):
            os.makedirs(directory + '/Output Folder/Pdf to Pptx')

        else: # If Folder already exists, then Skip/Pass
            pass
    except:
        pass

def createCacheFolder(directory): # Create Cache Folder for Pdf to Pptx Converter
    try:
        if not os.path.exists(directory + 'cache/Pdf to Pptx'):
            os.makedirs(directory + '/cache/Pdf to Pptx')

        else: # If Folder already exists, then Skip/Pass
            pass
    except:
        pass

def remove_pdf(directory): # Remove cache/unnecessary pptx files from directory
    files_in_directory = os.listdir(directory)
    filtered_files = [file for file in files_in_directory if file.endswith('.pdf')] # Filter format file (make sure program only remove .pptx files)

    for file in filtered_files:
        path_to_file = os.path.join(directory,file)
        os.remove(path_to_file)

# Main Program!!!
def pdftopptx_convert():
    directory = os.getcwd()

    createCacheFolder(directory)
    createFolder(directory)
    selectpdffile(directory)
    if n_files(directory) == 0:
        print("There is no files to convert")
        

    for filename in os.listdir(directory):
        if os.path.splitext(filename)[1] == '.pdf':
            print(Fore.BLUE+"Converting %s" % filename)
            print("Please wait a moment..."+Fore.WHITE)
            prs = Presentation()

            pages = convert_from_path(directory + '\\' + filename, 500,poppler_path = r"Resources\\Poppler\\bin")
            for index, page in enumerate(pages):
                #Save as 'jpg' in jpgs dir
                jpg_file = "cache/Pdf to Pptx/%s-(%d).jpg" % (filename,index)
                page.save(jpg_file, 'JPEG')

                #Get width/height of image
                image = Image.open(jpg_file)
                height = image.height
                width = image.width
                #Rotate 270 degrees if horizontal
                if height > width:
                    adjusted = image.rotate(270, expand=True)
                    adjusted.save(jpg_file)

                #Setup slide
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                left = top = 0
                slide.shapes.add_picture(jpg_file, left,top,height = prs.slide_height)

            image.close()
            prs.save('Output Folder/Pdf to Pptx/%s.pptx' % os.path.splitext(filename)[0])

    # Program Finished
    print(Fore.GREEN+"\nConversion Finished!"+Fore.WHITE)
    print("Clearing Cache...")
    time.sleep(2)
    remove_pdf(directory)
    converted_dir = directory+"\\Output Folder\\Pdf to Pptx"
    subprocess.Popen(f'explorer "{converted_dir}"')
    os.system("cls")
    cache_cleaner(directory)
